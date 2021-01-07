#!/usr/bin/env python
# coding: utf-8

# In[1]:


import PIL.ImageGrab
import time
import numpy as np
from PIL import Image, ImageChops
import os
from pptx import Presentation
from pptx.util import Inches
from natsort import natsorted
import comtypes.client


# In[2]:


number_frame = (2550,1480,2650,1520)
slide_frame = (1050,260,2700,1580)


# In[4]:


print('Insert maximal length of the Video (in minutes): ')
length = int(input())
print('Insert maximal amount of slides: ')
max_slides = int(input())


# In[5]:


entries = natsorted(os.listdir('Slides/'))
for entry in entries:
    name = 'Slides/' + entry
    os.remove(name)


# In[6]:


time.sleep(30)

start_time = time.time()

latest = PIL.ImageGrab.grab(slide_frame)
first_slide = PIL.ImageGrab.grab(slide_frame)
first_slide.save("Slides/slide1.png")


i = 2
while True:
    current_time = time.time()
    elapsed_time = current_time - start_time
    
    if(elapsed_time > (length*60)):
        print('Finished observation because of time!')
        break
    
    current = PIL.ImageGrab.grab(slide_frame)
    #current_array = np.array(current)
    #latest_array = np.array(latest)
    #diff = current_array - latest_array
    
    diff = ImageChops.difference(latest, current)
    diff = np.array(diff)
    
    if(diff.mean()>5):
        name = "Slides/slide" +str(i) +".png"
        i=i+1
        current.save(name)
        latest = current
    
    if(i>max_slides):
        print('Finished observation because of amount of slides!')
        break
    
    time.sleep(0.5)
    


# In[7]:


print('Deleting now Doublications!')
entries = natsorted(os.listdir('Slides/'))


# In[8]:


slides = []
for entry in entries:
    name = 'Slides/' + entry
    image = Image.open(name)
    unique = True

    if(len(slides)>0):
        for slide in slides:
            diff = ImageChops.difference(slide,image)
            diff = np.array(diff)
            
            if(diff.mean()<5):
                os.remove(name)
                unique = False
                print('Deleted picture: ' +name)
                break

        if(unique==True):
            slides.append(image)

    else:
        slides.append(image)


# In[9]:


print("Creating the PowerPoint now!")


# In[10]:


entries = natsorted(os.listdir('Slides/'))
prs = Presentation()


# In[11]:


for entry in entries:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left=Inches(0)
    top=Inches(0)
    name = 'Slides/' + entry
    img=slide.shapes.add_picture(name, Inches(0), Inches(0),
                               width= Inches(10), height= Inches(7))


# In[12]:


prs.save("AllSlides.pptx") # saving file


# In[13]:


def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()


# In[15]:


print('Procedure Finished!')


# In[ ]:




